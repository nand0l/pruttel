from pptx import Presentation

def create_presentation_from_md(md_file_path, pptx_template_path):
    # Read markdown file content
    with open(md_file_path, 'r', encoding='utf-8') as md_file:
        md_content = md_file.read()
    
    # Extract the title text (assuming it's the first line and starts with '# ')
    title_line = md_content.splitlines()[0]
    title_text = title_line.strip("# ").strip()
    
    # Split the title text
    parts = title_text.split("-")
    part1 = parts[0].strip()
    part2 = parts[1].strip()
    
    # Load the presentation template
    presentation = Presentation(pptx_template_path)
    
    # Use the first slide layout (assuming layout 0 is desired)
    slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(slide_layout)
    
    # Assuming placeholder indices are known (10 and 11 in this case)
    placeholder10 = slide.placeholders[10]  # Adjust index as needed
    placeholder10.text = part1
    
    placeholder11 = slide.placeholders[11]  # Adjust index as needed
    placeholder11.text = part2
    
    # Save the presentation
    output_pptx_path = md_file_path.replace('.md', '_presentation.pptx')
    presentation.save(output_pptx_path)
    print(f"Presentation saved to {output_pptx_path}")

# Example usage
md_file_path = 'C:\\code\\aws-bedrock\\pruttel\\streamlit\\output\\response_sonnet_1712044277.0.md'
pptx_template_path = 'C:\\code\\aws-bedrock\\pruttel\\streamlit\\pptx_source\\Skillsoft_GK_Template_empty.pptx'
create_presentation_from_md(md_file_path, pptx_template_path)

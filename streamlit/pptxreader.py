from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Load the presentation
pptx_file = 'Skillsoft_GK_Template_v01.2.pptx'
presentation = Presentation(pptx_file)

# Iterate through each slide layout available in the first slide master
for layout_index, layout in enumerate(presentation.slide_masters[0].slide_layouts):
    print(f"\nLayout {layout_index}: {layout.name}")
    # Iterate through placeholders in the layout
    for placeholder in layout.placeholders:
        # Check if placeholder is a shape
        if placeholder.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            print(f"  Placeholder {placeholder.placeholder_format.idx}: {placeholder.name}, Type: {placeholder.placeholder_format.type}")


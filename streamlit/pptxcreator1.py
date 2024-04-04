from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import re
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os


def read_markdown_content(md_file_path):
    """Reads and returns the content of a Markdown file."""
    with open(md_file_path, 'r', encoding='utf-8') as md_file:
        md_content = md_file.read()
    return md_content


def add_bullets_to_slide(slide, content):
    """Adds bullet points to a slide by creating a text box, using 'Spectral Light' font, and aligning text left."""
    # Define position and size of the text box
    left = Inches(1)  # 1 inch from the left edge of the slide
    top = Inches(1.5)  # 1.5 inches from the top of the slide
    width = Inches(8)  # Text box width
    height = Inches(5.5)  # Text box height

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    # Initialize the first paragraph
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT  # Align text to the left

    for line in content.split('\n'):
        if line.startswith('* ') or line.startswith('- '):
            # Start a new paragraph for a new bullet
            if p.text != '':  # if not the first bullet, add a new paragraph
                p = tf.add_paragraph()
            p.text = line.strip('* -')#[:5]
            p.level = 0
            p.font.size = Pt(28)
            p.font.name = "Spectral Light"
            p.font.bold = False
            #p.font.color.rgb = RGBColor(253, 239, 230)  # Hide first
            p.alignment = PP_ALIGN.LEFT

        # Detect sub-bullets with simple indentation
        elif line.startswith('  - ') or line.startswith('  + ') or line.startswith('  * '):
            p = tf.add_paragraph()
            p.text = line.strip('  +*')
            p.level = 1
            p.font.size = Pt(22)
            p.font.name = "Spectral Light"
            p.font.bold = False
            p.alignment = PP_ALIGN.LEFT
        # Detect sub-bullets with simple indentation
        elif line.startswith('    - ') or line.startswith('    + ') or line.startswith('    * '):
            p = tf.add_paragraph()
            p.text = line.strip('  +*')
            p.level = 2
            p.font.size = Pt(18)
            p.font.name = "Spectral Light"
            p.font.bold = False
            p.alignment = PP_ALIGN.LEFT
        # You can adjust the font size and level indentation as needed

# Remember to use this function within your markdown_to_presentation function when adding content to each slide.


def markdown_to_presentation(md_input, pptx_template_path, output_pptx_path):
    """Converts Markdown content to a PowerPoint presentation using a specified template."""
    presentation = Presentation(pptx_template_path)

    # Split the input into sections for title and subsequent slides
    sections = re.split(r'\n(?=## )', md_input)
    title_section = sections.pop(0)

    # Process title slide
    title_parts = title_section.split('\n', 1)[0].split('-')
    part1 = title_parts[0].strip("# ").strip()
    part2 = title_parts[1].strip()

    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    slide.placeholders[10].text = part1
    slide.placeholders[11].text = part2

    # Assuming layout 14 is for content slides
    slide_layout = presentation.slide_layouts[14]

    for section in sections:
        header, content = section.strip().split('\n', 1)
        slide_header = header.split(':', 1)[1].strip()

        slide = presentation.slides.add_slide(slide_layout)
        # Placeholder for slide header
        slide.placeholders[10].text = slide_header

        # Splitting content from notes
        content, _, notes_content = content.partition("_Notes:")

        # Add bullet points to the slide
        add_bullets_to_slide(slide, content)

        # Extract notes and add to the slide
        if notes_content:
            slide.notes_slide.notes_text_frame.text = notes_content.strip()
    
    # Assuming 'presentation' is your Presentation object and all other slides have been added
    thank_you_slide_layout = presentation.slide_layouts[19]  # Accessing the specific layout for the "Thank You" slide
    thank_you_slide = presentation.slides.add_slide(thank_you_slide_layout)

    # Setting text for specific placeholders
    thank_you_slide.placeholders[10].text = "Nando Lutgerink"
    thank_you_slide.placeholders[11].text = "nando.lutgerink@skillsoft.com"

    # Continue to save the presentation as before
    presentation.save(output_pptx_path)
    # Save the presentation
    presentation.save(output_pptx_path)
    print(f"Presentation saved to {output_pptx_path}")


# Paths for your Markdown file and PowerPoint template
md_file_path = 'C:\\code\\aws-bedrock\\pruttel\\streamlit\\output\\response_sonnet_1712062203.0.md'
filename_without_extension = os.path.splitext(
    os.path.basename(md_file_path))[0]
pptx_template_path = 'C:\\code\\aws-bedrock\\pruttel\\streamlit\\pptx_source\\Skillsoft_GK_Template_empty.pptx'
output_pptx_path = f'C:\\code\\aws-bedrock\\pruttel\\streamlit\\pptx_output\\{filename_without_extension}.pptx'

# Read Markdown content from file
md_content = read_markdown_content(md_file_path)

# Convert Markdown content to PowerPoint presentation
markdown_to_presentation(md_content, pptx_template_path, output_pptx_path)
